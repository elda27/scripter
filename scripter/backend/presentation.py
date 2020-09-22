from typing import List, Any
import pptx
from scripter.backend.presentation_base import PresentationBase
from scripter.backend.note_text import NoteText

class Presentation(PresentationBase):
    def __init__(self, filename):
        self.loaded_presentation = pptx.Presentation(filename)
        self.note_texts = []
        for p, slide in enumerate(self.loaded_presentation.slides):
            if slide.has_notes_slide:
                note = slide.notes_slide
                page_number = p + 1
                note_texts = [
                    paragraph.text
                    for paragraph in note.notes_text_frame.paragraphs
                ]
                self.note_texts.append(
                    NoteText(page_number, '\n'.join(note_texts))
                )

    def get_note_texts(self)->List[NoteText]:
        return self.note_texts

    def get_platform_presentation(self)->Any:
        return self.loaded_presentation
